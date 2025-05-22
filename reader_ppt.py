# # reader_ppt.py
# from pptx import Presentation

# def extract_text_ppt(ppt_path):
#     try:
#         prs = Presentation(ppt_path)
#         text = ""
#         for slide in prs.slides:
#             print(f"[INFO] Extracting text from slide {slide.slide_id}")
#             for shape in slide.shapes:
#                 if hasattr(shape, "text"):
#                     text += shape.text.strip() + " "
#         return text.replace("\n", " ").replace("\t", " ").strip()
#     except Exception as e:
#         print(f"[ERROR] PPT Extraction Failed: {e}")
#         return ""


from pptx import Presentation
from PIL import Image
import pytesseract
import os
import tempfile
from pdf2image import convert_from_path

def ocr_slide_as_image(ppt_path, slide_number):
    try:
        temp_pdf_path = os.path.join(tempfile.gettempdir(), "temp_ppt_slide.pdf")
        os.system(f'unoconv -f pdf -o "{temp_pdf_path}" "{ppt_path}"')

        if os.path.exists(temp_pdf_path):
            images = convert_from_path(temp_pdf_path)
            if slide_number < len(images):
                image = images[slide_number]
                return pytesseract.image_to_string(image).strip()
    except:
        return ""
    return ""

def extract_text_ppt(ppt_path):
    try:
        prs = Presentation(ppt_path)
        text = ""
        for idx, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text.strip() + " "
            if not slide_text.strip():
                ocr_text = ocr_slide_as_image(ppt_path, idx)
                if ocr_text:
                    slide_text += ocr_text + " "
            text += slide_text + " "
        return text.replace("\n", " ").replace("\t", " ").strip()
    except Exception as e:
        print(f"[ERROR] PPT Extraction Failed: {e}")
        return ""